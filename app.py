from flask import Flask, render_template, request, send_from_directory, Response, stream_with_context, session
import os
import uuid
from gtts import gTTS
from pptx import Presentation
import librosa
import soundfile as sf
from pydub import AudioSegment

app = Flask(__name__)
app.secret_key = "supersecretkey"  # necessario per le sessioni

# percorsi per i file caricati e l'audio finale
BASE_FOLDER = 'uploads'
app.config['UPLOAD_FOLDER'] = BASE_FOLDER
app.config['AUDIO_FOLDER'] = 'static/audio'

# crea le cartelle se non esistono
os.makedirs(BASE_FOLDER, exist_ok=True)
os.makedirs(app.config['AUDIO_FOLDER'], exist_ok=True)

# estrae testo dal .pptx
def estrai_testo_da_pptx(pptx_path):
    prs = Presentation(pptx_path)
    testo_slides = []
    for slide in prs.slides:
        testo_singola_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                testo_singola_slide.append(shape.text)
        testo_slides.append(" ".join(testo_singola_slide))
    return testo_slides

# trasforma il testo in audio (ma .mp3)
def trasforma_testo_in_audio(testo, output_audio):
    tts = gTTS(testo, lang="it")
    tts.save(output_audio)

# converte .mp3 in .wav usando librosa

def converti_da_mp3_a_wav(file_mp3):
    y, sr = librosa.load(file_mp3, sr=None)
    file_wav = file_mp3.replace(".mp3", ".wav")
    sf.write(file_wav, y, sr)
    os.remove(file_mp3)  # Rimuove il file .mp3 temporaneo
    return file_wav

# unisce file audio .wav
def unisci_file_audio(files_audio, output_audio):
    file_uniti = AudioSegment.empty()
    for file in files_audio:
        audio_segment = AudioSegment.from_wav(file)
        file_uniti += audio_segment
        os.remove(file)  # rimuove i file .wav temporanei
    file_uniti.export(output_audio, format="wav")

@app.route("/", methods=["GET", "POST"])
def index():
    if "session_id" not in session:
        session["session_id"] = str(uuid.uuid4())
    user_folder = os.path.join(BASE_FOLDER, session["session_id"])
    os.makedirs(user_folder, exist_ok=True)
    os.makedirs(app.config['AUDIO_FOLDER'], exist_ok=True)
    
    if request.method == "POST":
        pptx_file = request.files['pptx']
        pptx_path = os.path.join(user_folder, pptx_file.filename)
        pptx_file.save(pptx_path)
        
        testo_slides = estrai_testo_da_pptx(pptx_path)
        files_audio = []
        
        for i, testo_slide in enumerate(testo_slides):
            if testo_slide.strip():
                file_audio = os.path.join(app.config['AUDIO_FOLDER'], f"{session['session_id']}_slide_{i+1}.mp3")
                trasforma_testo_in_audio(testo_slide, file_audio)
                file_wav = converti_da_mp3_a_wav(file_audio)
                files_audio.append(file_wav)
        
        output_audio = os.path.join(app.config['AUDIO_FOLDER'], f"{session['session_id']}_audio_finale.wav")
        unisci_file_audio(files_audio, output_audio)

        os.remove(pptx_path)
        return render_template("index.html", audio_file=f"{session['session_id']}_audio_finale.wav")
    
    return render_template("index.html", audio_file=None)

@app.route("/uploads/<filename>")
def download_file(filename):
    return send_from_directory(app.config['AUDIO_FOLDER'], filename)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
