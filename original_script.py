import os
import librosa
import soundfile as sf
from gtts import gTTS
from pptx import Presentation
from pydub import AudioSegment

def estrai_testo_da_pptx(pptx_path):
    prs = Presentation(pptx_path)
    testo_slides = []
    for slide in prs.slides:
        testo_singola_slide = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                testo_singola_slide.append(shape.text)
        testo_slides.append(" ".join(testo_singola_slide))
    
    return testo_slides

def trasforma_testo_in_audio(testo, output_audio):
    tts = gTTS(testo, lang="it")
    tts.save(output_audio)

def converti_da_mp3_a_wav(file_mp3):
    file_wav = file_mp3.replace(".mp3", ".wav")
    y, sr = librosa.load(file_mp3, sr=None)
    sf.write(file_wav, y, sr)
    return file_wav

def unisci_file_audio(files_audio, output_audio):
    file_uniti = AudioSegment.empty()
    for file in files_audio:
        file_wav = converti_da_mp3_a_wav(file)
        audio_segment = AudioSegment.from_wav(file_wav)
        file_uniti += audio_segment
        os.remove(file_wav)
    file_uniti.export(output_audio.replace(".mp3", ".wav"), format="wav")

if __name__ == "__main__":
    
    pptx_file = "Lezione_3T.pptx"
    nome_file, estensione = pptx_file.split(".")
    output_audio = f"{nome_file}_AUDIO.mp3"

    testo_slides = estrai_testo_da_pptx(pptx_file)

    if testo_slides:
        files_audio = []

        for i, testo_slides in enumerate(testo_slides):
            if testo_slides.strip():
                file_audio = f"slide_{i+1}.mp3"
                print(f"Sto generando l'audio per la slide {i+1}...")
                trasforma_testo_in_audio(testo_slides, file_audio)
                files_audio.append(file_audio)
        
        print("\nSto unendo i file...")
        unisci_file_audio(files_audio, output_audio)

        for file in files_audio:
            os.remove(file)

        print(f"\nHo salvato il file audio finale con successo con il nome di: {output_audio}")

    else:
        print("Non ho trovato nessun testo nelle slide.")

    print("Ho terminato.")
